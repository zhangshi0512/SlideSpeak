from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import addphoto
import re
import pyttsx3


def presentate(defined_list, include_audio=True, include_code=True, include_images=True):
    prs = Presentation()
    engine = pyttsx3.init() if include_audio else None

    def add_slide(prs, layout, title, subtitle, generate_audio=False):
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title.strip()
        slide.placeholders[1].text = subtitle

        # Format title text
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(30)
            paragraph.font.bold = True
            paragraph.font.italic = False

        # Format subtitle text
        for paragraph in slide.placeholders[1].text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(16)
            paragraph.font.bold = False
            paragraph.font.italic = False

        if generate_audio and engine:
            full_text = f"{title}. {subtitle}"
            filename = title.replace(" ", "_") + ".mp3"
            engine.save_to_file(full_text, filename)
            engine.runAndWait()
        return slide

    def add_image_prompt_slide(prs, layout, prompt_text):
        slide = prs.slides.add_slide(layout)
        left = top = width = height = Inches(1)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = prompt_text
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        return slide

    # Define slide layouts
    text_slide_layout = prs.slide_layouts[1]
    image_slide_layout = prs.slide_layouts[6]

    # Clean and prepare summary text for each topic
    for d in defined_list:
        d["Summary"] = [re.sub(r'\d+\.\s+', '', item).strip()
                        for item in d["Summary"] if re.sub(r'\d+\.\s+', '', item).strip()]

    # Process each topic
    for item in defined_list:
        topic = item.get("Topic", "Untitled Topic")
        summary_lines = item.get("Summary", [])
        code = item.get("Code", "").strip()

        # Create summary slides
        if summary_lines:
            if len(summary_lines) > 5:  # Adjust threshold as needed
                mid = len(summary_lines) // 2
                summary_slide1 = "\n".join(summary_lines[:mid])
                summary_slide2 = "\n".join(summary_lines[mid:])
                add_slide(prs, text_slide_layout, topic,
                          summary_slide1, generate_audio=include_audio)
                add_slide(prs, text_slide_layout, topic,
                          summary_slide2, generate_audio=include_audio)
            else:
                add_slide(prs, text_slide_layout, topic, "\n".join(
                    summary_lines), generate_audio=include_audio)

        # Conditionally add the code snippet slide if valid code exists
        if include_code and code and "Error" not in code:
            add_slide(prs, text_slide_layout,
                      f"Code Snippet For {topic}", code, generate_audio=False)

        # Conditionally add image prompt slides
        if include_images:
            prompts = addphoto.get_images(topic, 2)
            for prompt in prompts:
                add_image_prompt_slide(prs, image_slide_layout, prompt)

    # Save the presentation
    prs.save("PPT.pptx")
