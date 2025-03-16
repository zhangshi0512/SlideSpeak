from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import win32com.client
import os
import re


### Helper method to preview the slide layouts available in a PowerPoint presentation
def previewSlideLayout():
    prs = Presentation()
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}:")
        for shape in layout.placeholders:
            print(f"  - Placeholder {shape.placeholder_format.idx}: {shape.name}")
        print("-" * 30)


def dictToPpt(inputDict: dict, speech_text=None):
    prs = Presentation() 

    # Add title slide
    addTitlePage(prs, inputDict["title"])

    # Process speech text into sections for speaker notes
    speech_sections = {}
    if speech_text:
        speech_sections = processSpecchIntoSections(speech_text, inputDict)

    # Add content slides
    for i, slide in enumerate(inputDict["slides"]):
        # Check if the slide has pre-generated speech
        slide_speech = slide.get("speech", "")
        
        # If no pre-generated speech, try to get it from speech_sections
        if not slide_speech and speech_text:
            slide_speech = speech_sections.get(slide["title"], "")
        
        # Add the slide with appropriate speaker notes
        addContentSlide(prs, slide, slide_speech)

    prs.save('PPT.pptx')

    print("Generation finished")

    # Uncomment if you want to use the shrinking functionality NOTE: Only works on Windows machine
    # shrinkTextInPowerpoint('PPT.pptx')


def processSpecchIntoSections(speech_text: str, inputDict: dict) -> dict:
    """
    Process the speech text and map sections to slide titles
    using explicit [SLIDE CHANGE] markers
    
    Args:
        speech_text: The entire speech text with [SLIDE CHANGE] markers
        inputDict: The presentation dictionary with slides
        
    Returns:
        A dictionary mapping slide titles to their corresponding speech sections
    """
    speech_sections = {}
    
    # Get all slide titles
    slide_titles = [slide["title"] for slide in inputDict["slides"]]
    
    # Split the speech by [SLIDE CHANGE] markers
    sections = speech_text.split("[SLIDE CHANGE]")
    
    # The first section is always for the first slide (introduction)
    if sections and slide_titles:
        speech_sections[slide_titles[0]] = sections[0].strip()
    
    # Map remaining sections to slides
    for i in range(1, min(len(sections), len(slide_titles))):
        speech_sections[slide_titles[i]] = sections[i].strip()
    
    # If we have more slides than sections, add empty notes for remaining slides
    for i in range(len(sections), len(slide_titles)):
        speech_sections[slide_titles[i]] = ""
    
    return speech_sections


### Helper method to add a title slide to the presentation, aka the first slide
### Take in the presentation object and the title: str of the presentation
def addTitlePage(prs: Presentation, title: str):
    # Presentation.slide_layouts[0] has the following layout (by default):
    #   - Placeholder 0: Title 1
    #   - Placeholder 1: Subtitle 2
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title.strip()
    slide.placeholders[1].text = "Subtitle Placeholder"
    setTitlePageStyle(slide)

## Helper method to set the style for the title slide, aka the first slide
def setTitlePageStyle(slide: Slide):
    if slide.shapes.title:
        title_shape = slide.shapes.title
        font = title_shape.text_frame.paragraphs[0].font
        font.size = Pt(44)
        font.bold = True


# Helper method, Add content of slide object inside slides from the inputDictp
# The input slideDict now supports the enhanced format from test.py:
# {
#   "title": str,
#   "content":[
#      {
#        "bulletPoint": str,
#        "shortSubPoints": [str, str, ...],
#        "details": [str, str, ...]
#      }, ... 
#   ],
#   "speech": str (optional)
# } 
def addContentSlide(prs: Presentation, slideDict: dict, speaker_notes=""):
    # Create the first slide for this topic
    currentSlide = addLayout1(prs, slideDict["title"])
    contentTextFrame = currentSlide.placeholders[1].text_frame
    wordCount = 0
    
    # First, add speaker notes to the first slide
    if hasattr(currentSlide, 'notes_slide'):
        notes_slide = currentSlide.notes_slide
        notes_textframe = notes_slide.notes_text_frame
        
        notes_content = []
        
        # Clean up the notes by removing special markers for TTS
        if speaker_notes:
            clean_notes = re.sub(r'\[PAUSE=\d+\]', '', speaker_notes)
            clean_notes = re.sub(r'\[SLIDE CHANGE\]', '', clean_notes)
            notes_content.append(clean_notes.strip())
        
        # If slide has speech directly from test.py format
        elif slideDict.get("speech"):
            clean_notes = re.sub(r'\[PAUSE=\d+\]', '', slideDict["speech"])
            clean_notes = re.sub(r'\[SLIDE CHANGE\]', '', clean_notes)
            notes_content.append(clean_notes.strip())
        
        # Now add the details from content
        details_text = []
        for content in slideDict["content"]:
            if content.get("details") and isinstance(content["details"], list) and content["details"]:
                bullet_point = content.get("bulletPoint", "").strip()
                details_text.append(f"\n\nDETAILS FOR: {bullet_point}")
                for detail in content["details"]:
                    details_text.append(f"- {detail.strip()}")
        
        if details_text:
            notes_content.append("\n".join(details_text))
        
        # Set the combined notes
        notes_textframe.text = "\n\n".join(notes_content)
    
    # Then add content, potentially creating additional slides if needed
    for content in slideDict["content"]:
        # The content now follows the enhanced structure from test.py
        # Extract the bulletPoint
        bulletPoint = content.get("bulletPoint", "")
        
        # Check if we have shortSubPoints, otherwise fall back to details
        shortSubPoints = content.get("shortSubPoints", [])
        details = content.get("details", [])
        
        # Use shortSubPoints if available, otherwise use details
        subPoints = shortSubPoints if shortSubPoints else details
        
        # Calculate word count for this content
        newContentWordCount = len(bulletPoint.split())
        if subPoints:
            newContentWordCount += sum([len(point.split()) for point in subPoints])
        
        # If content is too large, create a new slide
        if wordCount != 0 and wordCount + newContentWordCount > 120:
            currentSlide = addLayout1(prs, slideDict["title"])
            contentTextFrame = currentSlide.placeholders[1].text_frame
            wordCount = 0
            
            # Add details to speaker notes for additional slides too
            if hasattr(currentSlide, 'notes_slide'):
                notes_slide = currentSlide.notes_slide
                notes_textframe = notes_slide.notes_text_frame
                details_text = []
                
                # Only add details for content that will appear on this slide
                # (This is a simple approach; you might need to track which content items go on which slide)
                details_text.append(f"\nDETAILS FOR: {bulletPoint}")
                if content.get("details") and isinstance(content["details"], list):
                    for detail in content["details"]:
                        details_text.append(f"- {detail.strip()}")
                
                if details_text:
                    notes_textframe.text = "\n".join(details_text)
        
        p = contentTextFrame.add_paragraph()
        p.text = bulletPoint
        p.level = 0
        p.font.size = Pt(22)
        
        for point in subPoints:
            p = contentTextFrame.add_paragraph()
            point_text = point.strip()
            if point_text.startswith("- "):
                point_text = point_text[2:]
            p.text = point_text
            p.level = 1
            p.font.size = Pt(18)
        
        wordCount += newContentWordCount


def addLayout1(prs: Presentation, title: str) -> Slide:
    """ Add a slide with layout 1 (title and content) to the presentation. """
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title.strip()
    slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(40)
    slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    contentTextFrame = slide.placeholders[1].text_frame
    contentTextFrame.clear()
    contentTextFrame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    contentTextFrame.fit_text()
    # Remove the default empty paragraph
    if len(contentTextFrame.paragraphs) > 0:
        contentTextFrame.paragraphs[0]._element.getparent().remove(contentTextFrame.paragraphs[0]._element)

    return slide


def shrinkTextInPowerpoint(file_path: str) -> None:
    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = True

        presentation = powerpoint.Presentations.Open(os.path.abspath(file_path))

        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text_frame = shape.TextFrame
                    text_range = text_frame.TextRange
                    print(f"Found text frame with text : {text_range.Text}")

                    shrinkHappened = shrinkText(text_frame)
                    if shrinkHappened:
                        adjustFirstLineFontSize(text_frame)
        
        presentation.Save()
        presentation.Close()
        powerpoint.Quit()
    except Exception as e:
        print(f"Error: {e}")

def shrinkText(textFrame) -> bool:
    textRange = textFrame.TextRange
    shrunk = False
    while textRange.BoundHeight > textFrame.Parent.Height and textRange.Font.Size > 10:
        textRange.Font.Size -= 1
        shrunk = True
        print(f"Shrunk text to {textRange.Font.Size}")
    return shrunk

def adjustFirstLineFontSize(textFrame) -> None:
    textRange = textFrame.TextRange
    textContent = textRange.Text.strip()
    lines = textContent.split("\n")

    if lines:
        firstLine = lines[0]
        restText = "\n".join(lines[1:])
        textRange.Text = firstLine + "\n" + restText
        textRange.Words(1, len(firstLine)).Font.Size += 2

if __name__ == "__main__":
    import json
    # This is a test to demonstrate speaker notes functionality
    with open("./output/enriched_outline.json", "r", encoding="utf-8") as file:
        test_dict = json.load(file)
    
    with open("./output/presentation_speech.md", "r", encoding="utf-8") as file:
        test_speech = file.read()
    
    dictToPpt(test_dict, test_speech)