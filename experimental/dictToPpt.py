from pptx import Presentation
from pptx.slide import Slide
from pptx.util import Pt


### Helper method to preview the slide layouts available in a PowerPoint presentation
def previewSlideLayout():
    prs = Presentation()
    for i, layout in enumerate(prs.slide_layouts):
        print(f"Layout {i}:")
        for shape in layout.placeholders:
            print(f"  - Placeholder {shape.placeholder_format.idx}: {shape.name}")
        print("-" * 30)


testDict = {
    "title":"Coca-Cola's Use of Cloud Computing",
    "slides":[
        {
            "title":"Introduction",
            "content":[
                {
                    "bulletPoint":"Overview of Coca-Cola and its business model.",
                    "details":[
                        "Coca-Cola, founded in 1886 by John Pemberton, is a multinational beverage corporation headquartered in Atlanta, Georgia. It is one of the world's largest beverage companies, known for iconic products like Coca-Cola, Fanta, Sprite, and many others. Coca-Cola’s business model revolves around the production and distribution of soft drinks across various regions worldwide.",
                        "Coca-Cola operates through diverse channels including retail stores, vending machines, online platforms, and direct-to-consumer sales. The company also leverages strategic partnerships to expand its market reach."
                    ]
                },
                {
                    "bulletPoint":"Definition and importance of cloud computing in modern enterprises.",
                    "details":[
                        "Cloud computing refers to the delivery of computing services—such as servers, storage, databases, software, and analytics—over the Internet. These services are scalable and can be accessed from anywhere at any time.",
                        "Importance of Cloud Computing:",
                        "- Cost Efficiency: By shifting IT operations to cloud providers, Coca-Cola can reduce significant upfront costs associated with hardware and infrastructure. This allows for better allocation of resources towards core business activities.",
                        "- Scalability and Flexibility: Cloud computing enables quick scaling up or down based on demand, ensuring that the company’s technology infrastructure can adapt to varying levels of activity without downtime.",
                        "- Data Management: Cloud platforms offer robust data storage solutions that support data analytics, enabling Coca-Cola to make informed decisions based on real-time insights from consumer behavior and market trends."
                    ]
                }
            ]
        },
        {
            "title":"Background on Cloud Computing",
            "content":[
                {
                    "bulletPoint":"Key benefits: Scalability, cost reduction, flexibility, data security.",
                    "details":[
                        "Scalability allows businesses to expand or contract their cloud resources according to demand, ensuring that they can handle fluctuations in traffic and sales without significant upfront investment.",
                        "Cost reduction is achieved through pay-as-you-go models, which eliminate the need for expensive on-premises infrastructure, thereby saving both initial capital expenses and ongoing operational costs.",
                        "Flexibility enables companies to quickly deploy new services or applications by leveraging pre-configured cloud environments, which can be scaled up or down as needed. This agility is particularly valuable in today's fast-paced business environment.",
                        "Data security measures include advanced encryption techniques, multi-factor authentication, and regular audits to protect sensitive information from unauthorized access and breaches."
                    ]
                },
                {
                    "bulletPoint":"Common cloud services: Infrastructure as a Service (IaaS), Platform as a Service (PaaS), Software as a Service (SaaS).",
                    "details":[
                        "Infrastructure as a Service (IaaS) provides businesses with virtualized computing resources, including servers, storage, and networking. Companies like Coca-Cola can use IaaS to host their applications or databases without the need for physical hardware.",
                        "Platform as a Service (PaaS) offers development tools and services that allow developers to build, test, deploy, and manage software applications. This service enables companies to focus on application development rather than infrastructure management.",
                        "Software as a Service (SaaS) provides complete software solutions over the internet, eliminating the need for local installation and maintenance. Examples include customer relationship management (CRM) systems, human resource management (HRM) tools, and accounting software."
                    ]
                }
            ]
        },
        {
            "title":"Coca-Cola's Cloud Strategy",
            "content":[
                {
                    "bulletPoint":"Objectives of Coca-Cola’s cloud adoption.",
                    "details":[
                        "Cost reduction and efficiency improvements through cloud computing: Migrating to the cloud helps reduce on-premises IT infrastructure costs, including hardware, software, and maintenance expenses. This allows the company to reinvest savings into other business areas.",
                        "Enhanced scalability and flexibility: Cloud computing enables Coca-Cola to scale resources up or down based on demand without significant upfront investments, allowing for better alignment with fluctuating business needs.",
                        "Data analysis and insights: Utilizing cloud-based analytics tools to gain deeper insights from customer data, sales performance, and market trends. For example, analyzing consumer behavior patterns to inform marketing campaigns more effectively.",
                        "Global collaboration and communication: Cloud platforms support seamless real-time communication and collaboration among global teams, improving overall efficiency and productivity."
                    ]
                },
                {
                    "bulletPoint":"Decision-making process: Internal factors, external influences.",
                    "details":[
                        "Internal factors influencing the decision: Current IT infrastructure, existing software applications, employee skill sets, and internal processes. For instance, assessing whether current in-house systems can effectively support cloud migration without extensive retooling.",
                        "External influences on the decision: Market trends, industry competitors’ strategies, regulatory requirements, and technological advancements. Evaluating how other companies in the beverage industry are adopting cloud technologies to stay competitive and compliant with data privacy laws.",
                        "Risk assessment considerations: Security risks associated with cloud storage, potential vendor lock-in scenarios, and compliance challenges related to data protection regulations. Implementing robust security measures and considering multiple vendors to mitigate these risks."
                    ]
                }
            ]
        },
        {
            "title":"Key Applications in Operations",
            "content":[
                {
                    "bulletPoint":"Supply chain management systems - real-time inventory tracking.",
                    "details":[
                        "Coca-Cola leverages cloud-based supply chain management solutions to monitor and optimize its global inventory. This allows for real-time updates on stock levels, reducing the risk of stockouts or overstocking.",
                        "Examples include using IoT devices and sensors to track temperatures and locations of goods in transit, ensuring product quality and compliance with safety standards."
                    ]
                },
                {
                    "bulletPoint":"Customer relationship management (CRM) tools - enhanced customer service and analytics.",
                    "details":[
                        "Coca-Cola utilizes cloud CRM systems to collect and analyze customer data, improving the efficiency and effectiveness of their customer interactions. This includes tracking consumer preferences, purchase history, and feedback through various channels like social media, websites, and mobile apps.",
                        "CRM tools help in personalizing marketing communications and offering tailored products or services based on individual customer profiles."
                    ]
                },
                {
                    "bulletPoint":"Marketing automation platforms - targeted advertising campaigns.",
                    "details":[
                        "Coca-Cola employs cloud-based marketing automation software to create, execute, and analyze targeted advertising campaigns. These tools enable the company to automate email marketing, social media ads, and other digital marketing activities.",
                        "By integrating CRM data with marketing automation platforms, Coca-Cola can segment its customer base into various groups based on demographics, behaviors, or purchase history, allowing for more personalized and effective marketing strategies."
                    ]
                }
            ]
        },
        {
            "title":"Technical Implementation",
            "content":[
                {
                    "bulletPoint":"Selection of cloud providers: Amazon Web Services, Microsoft Azure, Google Cloud Platform.",
                    "details":[
                        "Coca-Cola has chosen to leverage leading cloud service providers such as Amazon Web Services (AWS), Microsoft Azure, and Google Cloud Platform (GCP) for their robust infrastructure capabilities and global presence.",
                        "These providers offer a wide range of services including storage solutions, database management, machine learning tools, and analytics platforms which align with Coca-Cola’s diverse operational needs."
                    ]
                },
                {
                    "bulletPoint":"Security measures in place: Compliance with regulations like GDPR, data encryption.",
                    "details":[
                        "Coca-Cola has implemented stringent security measures to protect sensitive company information. This includes ensuring compliance with international data protection regulations such as the General Data Protection Regulation (GDPR), which mandates strict controls over how personal data is processed and stored.",
                        "Data Encryption: All critical data is encrypted both in transit and at rest to prevent unauthorized access. This ensures that even if data were to be intercepted, it would not be readable without the appropriate decryption key.",
                        "Access Controls: Multi-factor authentication (MFA) and role-based access controls are enforced to ensure that only authorized personnel have access to specific resources or information."
                    ]
                }
            ]
        },
        {
            "title":"Benefits and Challenges",
            "content":[
                {
                    "bulletPoint":"Operational efficiencies - reduced downtime, improved collaboration.",
                    "details":[
                        "Cloud computing enables Coca-Cola to reduce downtime by automatically scaling resources up or down based on demand. This ensures the company can maintain high performance during peak usage times without overprovisioning infrastructure.",
                        "Collaboration is enhanced through cloud-based tools that allow employees and teams across different locations to access data and applications from anywhere, promoting real-time communication and project management."
                    ]
                },
                {
                    "bulletPoint":"Financial implications - cost savings on hardware, maintenance.",
                    "details":[
                        "Cloud computing helps Coca-Cola save significantly on the costs associated with purchasing and maintaining physical servers. By moving to a cloud infrastructure, the company can avoid large upfront capital expenditures for hardware and reduce ongoing operational expenses.",
                        "Maintenance is outsourced to the cloud service provider, eliminating the need for in-house IT teams to handle software updates, security patches, and other technical maintenance tasks, which further reduces costs and improves reliability."
                    ]
                }
            ]
        },
        {
            "title":"Case Studies and Success Stories",
            "content":[
                {
                    "bulletPoint":"Real-world examples of cloud initiatives at Coca-Cola.",
                    "details":[
                        "Coca-Cola implemented a cloud-based inventory management system to streamline supply chain operations, reducing stockouts and excess inventory by 25%. This initiative allowed for more efficient distribution and reduced costs associated with logistics by $10 million annually.",
                        "The company also adopted cloud computing for its marketing analytics platform. By leveraging big data from social media and customer feedback, Coca-Cola was able to personalize marketing campaigns, increasing engagement rates by 30% across various platforms such as Facebook and Instagram."
                    ]
                },
                {
                    "bulletPoint":"Impact on business performance: Quantifiable results and improvements.",
                    "details":[
                        "Cloud computing enabled Coca-Cola to improve its customer service operations through a real-time monitoring system. This resulted in a 40% reduction in response times for customer inquiries, leading to higher customer satisfaction scores.",
                        "In terms of cost savings, cloud infrastructure allowed the company to reduce IT maintenance costs by 35%. Additionally, by adopting cloud services, Coca-Cola was able to allocate more resources to innovation and product development, enhancing its competitive edge.",
                        "Operational efficiency also saw a significant boost as cloud solutions facilitated better collaboration across global teams. For instance, project management tools on the cloud enabled faster decision-making processes, which translated into a 20% increase in productivity among remote workers."
                    ]
                }
            ]
        },
        {
            "title":"Future Outlook",
            "content":[
                {
                    "bulletPoint":"Projected growth in cloud usage across the organization.",
                    "details":[
                        "Coca-Cola anticipates a significant increase in cloud-based operations, leveraging services like AWS and Microsoft Azure to streamline processes and reduce costs.",
                        "The company plans to migrate more applications and data storage onto the cloud within the next three years, aiming for over 90% of its IT infrastructure to be cloud-hosted by 2025.",
                        "Cloud computing will enable real-time analytics and decision-making across various departments such as marketing, sales, and supply chain management."
                    ]
                },
                {
                    "bulletPoint":"Potential integration with IoT technologies for enhanced data collection.",
                    "details":[
                        "IoT devices can be integrated into the cloud to collect real-time data on consumer behavior, store conditions, and equipment performance, providing valuable insights for product development and supply chain optimization.",
                        "For example, sensors in Coca-Cola vending machines could monitor temperature and inventory levels, sending this data directly to the cloud for analysis and automated restocking alerts.",
                        "Integration with smart packaging technologies could offer personalized marketing messages based on consumer preferences and behaviors, enhancing customer engagement and loyalty."
                    ]
                }
            ]
        },
        {
            "title":"Conclusion",
            "content":[
                {
                    "bulletPoint":"Summary of key points discussed.",
                    "details":[
                        "Coca-Cola has leveraged cloud computing to streamline its operations, enhance supply chain efficiency, and improve data analytics capabilities. This includes the use of cloud-based tools for inventory management, order fulfillment, and real-time tracking of products across multiple channels.",
                        "The implementation of advanced analytics on the cloud has enabled Coca-Cola to gain deeper insights into consumer behavior, preferences, and market trends. These insights drive personalized marketing campaigns and product development initiatives that better align with customer needs.",
                        "Integration of cloud solutions has also facilitated collaboration among different departments within the company. For instance, sales teams can access real-time data on inventory levels, while marketing teams have instant access to customer feedback and analytics reports."
                    ]
                },
                {
                    "bulletPoint":"Impact on Coca-Cola's competitive edge and future strategies.",
                    "details":[
                        "By adopting cloud computing, Coca-Cola has gained a significant competitive advantage in terms of operational efficiency, agility, and data-driven decision making. This allows the company to respond faster to market changes and customer demands.",
                        "Future strategies may include expanding use cases for cloud technologies such as IoT (Internet of Things) integration for smart logistics, AI-powered predictive maintenance, and enhanced cybersecurity measures to protect sensitive data.",
                        "Furthermore, continued investment in cloud infrastructure will support innovation in areas like personalized marketing solutions, augmented reality experiences, and sustainable supply chain practices."
                    ]
                }
            ]
        }
    ]
}


def dictToPpt(inputDict: dict):
    prs = Presentation() 

    # Add title slide
    addTitlePage(prs, inputDict["title"])


    prs.save('PPT.pptx')




### Helper method to add a title slide to the presentation
def addTitlePage(prs: Presentation, title: str):
    ### Take in the presentation object and the title: str of the presentation
    # Presentation.slide_layouts[0] has the following layout:
    # Lyout 0:
    #   - Placeholder 0: Title 1
    #   - Placeholder 1: Subtitle 2
    #   - Placeholder 10: Date Placeholder 3
    #   - Placeholder 11: Footer Placeholder 4
    #   - Placeholder 12: Slide Number Placeholder 5
    titleSlideLayout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(titleSlideLayout)

    slide.shapes.title.text = title.strip()
    slide.placeholders[1].text = "Subtitle Placeholder"
    

def setTitlePageStyle(slide: Slide):
    if slide.shapes.title:
        title_shape = slide.shapes.title
        font = title_shape.text_frame.paragraphs[0].font
        font.size = Pt(44)
        font.bold = True
        font.italic = True

dictToPpt(testDict)


